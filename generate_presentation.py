from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Google Cloud Design System Colors
COLOR_BG = RGBColor(255, 255, 255)        # White
COLOR_TEXT_MAIN = RGBColor(60, 64, 67)    # #3C4043 (Dark Charcoal)
COLOR_TEXT_LIGHT = RGBColor(112, 117, 122) # #70757A (Medium Gray)
COLOR_BLUE = RGBColor(66, 133, 244)       # #4285F4
COLOR_RED = RGBColor(234, 67, 53)         # #EA4335
COLOR_YELLOW = RGBColor(251, 188, 4)      # #FBBC04
COLOR_GREEN = RGBColor(52, 168, 83)       # #34A853
COLOR_GRAY_LIGHT = RGBColor(241, 243, 244) # #F1F3F4

# Fonts
FONT_HEADING = "Arial" 
FONT_BODY = "Arial"

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_google_pattern(slide):
    # Add a subtle geometric pattern on the right
    # Blue Triangle
    left = Inches(8.5)
    top = Inches(0.5)
    width = Inches(1)
    height = Inches(1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BLUE
    shape.line.fill.background()
    shape.rotation = 90

    # Red Triangle
    left = Inches(9.2)
    top = Inches(1.2)
    width = Inches(0.8)
    height = Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_RED
    shape.line.fill.background()
    shape.rotation = 270

    # Yellow Triangle
    left = Inches(8.8)
    top = Inches(1.8)
    width = Inches(0.6)
    height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_YELLOW
    shape.line.fill.background()

    # Green Triangle
    left = Inches(9.5)
    top = Inches(2.2)
    width = Inches(0.5)
    height = Inches(0.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_GREEN
    shape.line.fill.background()
    shape.rotation = 180

def add_footer(slide):
    # "Proprietary + Confidential" top right
    left = Inches(7)
    top = Inches(0.2)
    width = Inches(2.5)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Proprietary + Confidential"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.RIGHT

    # "Google Cloud" bottom right
    left = Inches(8)
    top = Inches(7.0)
    width = Inches(1.5)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Google Cloud"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT

def create_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_BG)
    
    # Google Cloud Logo (Text representation) top left
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(2)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Google Cloud"
    p.font.name = FONT_HEADING
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.font.bold = True

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(6.5)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(40)
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    top = Inches(4.5)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.alignment = PP_ALIGN.LEFT

    add_google_pattern(slide)
    
    # Date bottom left
    top = Inches(6.8)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "February 2026"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_LIGHT

    return slide

def create_content_slide(prs, title_text, content_items):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_BG)
    
    add_footer(slide)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_BLUE
    p.font.bold = True
    
    # Content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(12)
        if item.startswith("  "):
            p.level = 1
        
    return slide

def create_presentation():
    prs = Presentation()
    
    # --- Slide 1: Title Slide ---
    slide1 = create_title_slide(
        prs, 
        "Building the Software Stack 'Community Engine'", 
        "Hustle & Have Fun: Redefining UX/PM/Eng Collaboration\n\nPresenters: [Your Name] (UX), Brittany (PM), Rob (Eng)"
    )
    slide1.notes_slide.notes_text_frame.text = "Hi everyone. Today, Brittany, Rob, and I want to share a story about how we completely reimagined how UX, Product, and Engineering collaborate. By leaning into the principles of 'Stay Scrappy' and 'Hustle & Have Fun', we moved from a blank whiteboard to a fully functional 'Community Engine' MVP in exactly 3 months."

    # --- Slide 2: The Opportunity & The "Scrappy" Shift ---
    content2 = [
        "User Problem: External contributors faced a 'Black Box' when trying to build on our TPU software stack—it was too hard to know what was broken or where to start.",
        "Business Problem: Internal engineers were drowning in review bottlenecks and bespoke support requests, slowing down our own feature velocity.",
        "The Opportunity (Open Acceleration): Enable the community to drive performance enhancements (like vLLM) so internal engineers don't have to.",
        "The Mindset ('Stay Scrappy'): Stop standard, pixel-perfect design cycles. We rapidly brainstormed the 'Six-Stage Flywheel' concepts to find our hook, and then immediately jumped into building a working MVP (the end-to-end automated flow)."
    ]
    slide2 = create_content_slide(prs, "The Opportunity & The 'Scrappy' Shift", content2)
    slide2.notes_slide.notes_text_frame.text = "The core issue was that community developers wanted to build on our TPU software stack, but faced a 'Black Box'. Internally, our engineers were drowning in bespoke requests and review bottlenecks. Instead of doing what we normally do—spending weeks on high-fidelity Figma mockups and PM requirements docs—we got scrappy. We defined a 'Six-Stage Flywheel' for community engagement, and immediately started writing code to build the foundational automation MVP."

    # --- Slide 3: The MVP – The "Community Engine" ---
    content3 = [
        "What we built: A 3-step automated workflow requiring zero human intervention:",
        "  1. Auto CUJ (The Truth Pipeline): Automated sweeps find TPU test failures and expose them in a transparent matrix.",
        "  2. Auto Mission Board (Candy Crush): Those failures automatically generate 'Good First Issues' to lower the entry barrier for contributors.",
        "  3. Auto Leaderboard: Triggers an automated 'Inclusive Contributor Wall' update when code is merged, rewarding the dopamine loop.",
        "\n[INSERT SCREENSHOT OF MISSION BOARD / CONTRIBUTOR WALL HERE]"
    ]
    slide3 = create_content_slide(prs, "The MVP – The 'Community Engine'", content3)
    slide3.notes_slide.notes_text_frame.text = "Our MVP focused purely on the automated workflow to eliminate human overhead. We call it the Community Engine. Step 1 is the Auto CUJ, which automatically detects missing features and surfaces them. Step 2 is the 'Candy Crush' Mission Board: those gaps automatically turn into bite-sized GitHub issues. Step 3 closes the loop—when the community merges a fix, an automated Leaderboard updates to celebrate them. We went from concept to a live, automated ecosystem in 3 months."

    # --- Slide 4: Roles Reimagined ("Mission First") ---
    content4 = [
        "UX / Design ([Your Name]): Erased the 'mockup' boundary. Wrote the actual working Python code and GitHub Actions to automate the MVP.",
        "Engineering (Rob): Shifted to an accelerator. Helped modularize the kernel tuner, solved complex infrastructure blocks, and reviewed the designer's code.",
        "Product (Brittany): Managed the 'Hockey Stick'. Aligned the internal strategy, prioritized the 'Golden Path', and provided continuous input without stalling execution."
    ]
    slide4 = create_content_slide(prs, "Roles Reimagined ('Mission First')", content4)
    slide4.notes_slide.notes_text_frame.text = "The secret to doing this in 3 months was erasing our traditional job titles and working 'Mission First'. I stepped outside UX and actually wrote the working Python scripts and GitHub Actions to bring the Mission Board to life. Rob shifted from strictly writing features to acting as an accelerator—reviewing my code and modularizing the internal tuner. And Brittany focused on unblocking us and aligning our work with the wider TPU strategy without slowing down our daily execution."

    # --- Slide 5: The Impact & The Future ---
    content5 = [
        "Impact: A live, scalable Contribution Ecosystem that reduces internal engineering load and increases external velocity.",
        "Key Takeaway: Proving value via working code (Stay Scrappy) builds momentum faster than pitching slides.",
        "Our Ask: How can leadership help us scale this embedded, high-velocity ('Hustle & Have Fun') team model to other critical infrastructure tracks?"
    ]
    slide5 = create_content_slide(prs, "The Impact & The Future", content5)
    slide5.notes_slide.notes_text_frame.text = "The impact here isn't just a cool GitHub script. We've built a scalable engine that actively reduces our internal engineering load. What we learned is that proving value via working code builds unstoppable momentum. Our ask for you today is this: How can we scale this exact model—this high-velocity, boundary-less collaboration—to other critical infrastructure teams at Google?"

    # Save presentation
    output_filename = "Community_Engine_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation generated successfully: {output_filename}")

if __name__ == "__main__":
    create_presentation()
